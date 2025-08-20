import pptx
import cv2
from lxml import etree
import inspect
import subprocess
import scene
import os

ppt_name = "Defense.pptx"
movie_directory = os.getcwd() + "/media/videos/scene/720p15/sections/"


def getFirstFrame(videofile, name):
    vidcap = cv2.VideoCapture(videofile + ".mp4")
    success, image = vidcap.read()
    if success:
        cv2.imwrite(f"first_frames/{name}.png", image)


def xpath(el, query):
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    return etree.ElementBase.xpath(el, query, namespaces=nsmap)


def autoplay_media(media):
    el_id = xpath(media.element, './/p:cNvPr')[0].attrib['id']
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), './/p:cond')[0]
    cond.set('delay', '0')


scenes = {
    "TitleScene": 2,
    "AstroScene": 9,
    "BHScene": 3,
    "SystemIntroScene": 5,
    "ProblemStatement": 9,
    "VariableScene": 6,
    "HamiltonianScene": 11,
    "NoTheta": 2,
    "ZVCScene": 7,
    "CircularOrbitScene": 9,
    "RescaleScene": 9,
    "SymmetryScene": 6,
    "EscapeDefinitionScene": 15,
    "EscapeQuantities": 4,
    "EscapePlot": 10,
    "EscapePlot1": 2,
    "EscapePlot2": 6,
    "TrapDemo": 3,
    "ChaosDefinitionScene": 7,
    "SaliScene": 17,
    "SaliPlot": 4,
    "SaliTrap": 2,
    "SaliTrapPlot": 5,
    "CountOrdered": 6,
    "MBHamiltonianScene": 13,
    "MBHZVC": 6,
    "MBHRescale": 6,
    "MBHEscapePlot": 10,
    "MBHChaos": 3,
    "MBHConclusion": 7,
    "Conclusion": 5
}

def generate_slides():
    ppt = pptx.Presentation()
    ppt.slide_width = 9144000
    ppt.slide_height = 5143500
    for scene, section_number in scenes.items():
        for i in [f"{scene}_{j:04}_autocreated" if j == 0 else f"{scene}_{j:04}_unnamed" for j in range(section_number)]:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            movie_file = f"{movie_directory}/{i}"
            
            getFirstFrame(movie_file, i)
            media = slide.shapes.add_movie(
                f"{movie_file}.mp4", 0, 0, ppt.slide_width, ppt.slide_height,
                poster_frame_image=f"first_frames/{i}.png"
            )
            autoplay_media(media)
    ppt.save(ppt_name)

def generate_scenes():
    classes = [name for name, obj in inspect.getmembers(scene, inspect.isclass) if obj.__module__ == scene.__name__]
    for i in classes:
        subprocess.run(["manim", "-qm", "scene.py", i, "--fps", "15", "--save_sections"])

def main():
    generate_slides()

if __name__ == '__main__':
    main()
